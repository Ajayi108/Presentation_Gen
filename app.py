import io
import json
import os
import re
import urllib.error
import urllib.parse
import urllib.request
from pathlib import Path
from typing import Any

import streamlit as st

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - handled in UI
    Presentation = None
    RGBColor = None
    Inches = None
    Pt = None

APP_NAME = "ai_presentation_generator"
DEFAULT_MODEL = "gemini-2.5-flash"
API_URL = "https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
UNSPLASH_RANDOM_URL = "https://api.unsplash.com/photos/random"
OUTPUT_DIR = Path("generated_presentations")

THEME_PRESETS = {
    "Editorial Sand": {
        "primary": "#2F4858",
        "secondary": "#F7F1E8",
        "accent": "#D97B4E",
        "surface": "#FFF9F2",
        "ink": "#1E252B",
        "muted": "#5A676F",
        "description": "Warm and polished for product demos, consulting decks, and strategy work.",
    },
    "Neo Mint": {
        "primary": "#173C3A",
        "secondary": "#ECF8F5",
        "accent": "#3CBFA1",
        "surface": "#F7FFFD",
        "ink": "#122221",
        "muted": "#4A6661",
        "description": "Fresh, modern, and clean for startups, AI, and innovation topics.",
    },
    "Midnight Luxe": {
        "primary": "#171A2D",
        "secondary": "#F4F1EA",
        "accent": "#C88C3A",
        "surface": "#FCFAF6",
        "ink": "#181A1F",
        "muted": "#5C5E68",
        "description": "High-contrast and premium for investor pitches and executive briefings.",
    },
    "Sage Minimal": {
        "primary": "#32463D",
        "secondary": "#F2F5EF",
        "accent": "#90A955",
        "surface": "#FAFCF8",
        "ink": "#1F2824",
        "muted": "#5E6D66",
        "description": "Soft and restrained for education, sustainability, and nonprofit decks.",
    },
    "Coral Studio": {
        "primary": "#7A2E2E",
        "secondary": "#FFF1EC",
        "accent": "#FF7A59",
        "surface": "#FFF8F5",
        "ink": "#2F1D1D",
        "muted": "#7D5B58",
        "description": "Energetic and warm for marketing launches and creative decks.",
    },
    "Ocean Signal": {
        "primary": "#10324A",
        "secondary": "#EAF5FB",
        "accent": "#2FA7D8",
        "surface": "#F6FBFE",
        "ink": "#14232C",
        "muted": "#5E7482",
        "description": "Clear and technical for analytics, research, and product strategy.",
    },
    "Lavender Grid": {
        "primary": "#443C68",
        "secondary": "#F2F0FB",
        "accent": "#8B7CF6",
        "surface": "#FAF9FE",
        "ink": "#201D2F",
        "muted": "#66627A",
        "description": "Contemporary and elegant for design-forward storytelling.",
    },
    "Amber Slate": {
        "primary": "#393939",
        "secondary": "#F7F3EA",
        "accent": "#D9A441",
        "surface": "#FCFAF4",
        "ink": "#222120",
        "muted": "#676057",
        "description": "Balanced and mature for operations, finance, and advisory decks.",
    },
    "Rose Quartz": {
        "primary": "#6A3E52",
        "secondary": "#FAEEF3",
        "accent": "#E47BAA",
        "surface": "#FFF8FB",
        "ink": "#2B1F25",
        "muted": "#7B6670",
        "description": "Soft premium styling for fashion, wellness, and lifestyle presentations.",
    },
    "Forest Signal": {
        "primary": "#18392B",
        "secondary": "#EDF6F1",
        "accent": "#2E8B57",
        "surface": "#F7FCF9",
        "ink": "#16211C",
        "muted": "#5B6E63",
        "description": "Grounded and credible for climate, impact, and sustainability work.",
    },
    "Cobalt Pulse": {
        "primary": "#203A8F",
        "secondary": "#EEF2FF",
        "accent": "#4F7CFF",
        "surface": "#F8FAFF",
        "ink": "#182032",
        "muted": "#66718C",
        "description": "Sharp and energetic for software, AI, and technical demos.",
    },
    "Terracotta Paper": {
        "primary": "#6B3F2C",
        "secondary": "#F8EFE9",
        "accent": "#C86B42",
        "surface": "#FFF9F6",
        "ink": "#2C221F",
        "muted": "#76615A",
        "description": "Textured editorial warmth for brand, narrative, and story-led decks.",
    },
}

IMAGE_LAYOUT_OPTIONS = {
    "Right aligned throughout": "right",
    "Left aligned throughout": "left",
    "Alternating left and right": "alternate",
}

PRESENTATION_SCHEMA = {
    "type": "OBJECT",
    "properties": {
        "title": {"type": "STRING"},
        "subtitle": {"type": "STRING"},
        "slides": {
            "type": "ARRAY",
            "items": {
                "type": "OBJECT",
                "properties": {
                    "title": {"type": "STRING"},
                    "bullets": {"type": "ARRAY", "items": {"type": "STRING"}},
                    "speaker_notes": {"type": "STRING"},
                },
                "required": ["title", "bullets", "speaker_notes"],
            },
        },
        "closing_message": {"type": "STRING"},
    },
    "required": ["title", "subtitle", "slides", "closing_message"],
}

st.set_page_config(page_title="DeckMuse", page_icon=":material/slideshow:", layout="wide")

@st.cache_data(show_spinner=False)
def load_dotenv(path: str = ".env") -> dict[str, str]:
    env_path = Path(path)
    loaded: dict[str, str] = {}
    if not env_path.exists():
        return loaded
    for raw_line in env_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        os.environ.setdefault(key, value)
        loaded[key] = value
    return loaded

def get_api_key() -> str | None:
    load_dotenv()
    return os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")

def get_unsplash_key() -> str | None:
    load_dotenv()
    return os.getenv("UNSPLASH_ACCESS_KEY")

def parse_project_brief(raw_text: str) -> list[str]:
    # Accept simple bullet-style input and keep only the first five lines.
    bullets: list[str] = []
    for line in raw_text.splitlines():
        cleaned = re.sub(r"^[\-\*\u2022\s]+", "", line.strip())
        if cleaned:
            bullets.append(cleaned)
    return bullets[:5]

def build_prompt(topic: str, audience: str, tone: str, slide_count: int, theme_name: str, project_brief: list[str]) -> str:
    # Feed the optional project brief into Gemini so the deck reflects the full context.
    brief_block = ""
    if project_brief:
        brief_lines = "\n".join(f"- {item}" for item in project_brief)
        brief_block = f"\nProject description bullets:\n{brief_lines}\n"
    return f"""
Create a professional presentation outline as JSON.

Topic: {topic}
Audience: {audience}
Tone: {tone}
Visual direction: {theme_name}
Requested content slides: {slide_count}
{brief_block}
Instructions:
- Return valid JSON only.
- Create a strong presentation title and a one-sentence subtitle.
- Provide exactly {slide_count} content slides in the slides array.
- Each slide needs a concise title, 3 to 4 bullets, and speaker_notes.
- Use the optional project description bullets when they are provided.
- Bullets should be short, specific, and presentation-ready.
- closing_message should be a short ending line for the final slide.
""".strip()

def call_gemini(prompt: str, api_key: str, model: str) -> dict[str, Any]:
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.7,
            "responseMimeType": "application/json",
            "responseSchema": PRESENTATION_SCHEMA,
        },
    }
    request = urllib.request.Request(
        API_URL.format(model=model),
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json", "X-goog-api-key": api_key},
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=60) as response:
            raw = response.read().decode("utf-8")
    except urllib.error.HTTPError as exc:
        details = exc.read().decode("utf-8", errors="ignore")
        raise RuntimeError(f"Gemini API request failed ({exc.code}). {details}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"Network error while contacting Gemini: {exc.reason}") from exc
    response_json = json.loads(raw)
    parts = response_json.get("candidates", [{}])[0].get("content", {}).get("parts", [])
    text = "".join(part.get("text", "") for part in parts).strip()
    if not text:
        raise RuntimeError("Gemini returned an empty response.")
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, re.S)
        if not match:
            raise RuntimeError("Gemini did not return valid JSON.")
        return json.loads(match.group(0))

def with_referral(url: str | None) -> str:
    if not url:
        return ""
    separator = "&" if "?" in url else "?"
    return f"{url}{separator}utm_source={APP_NAME}&utm_medium=referral"

def fetch_json(url: str, headers: dict[str, str]) -> dict[str, Any]:
    request = urllib.request.Request(url, headers=headers, method="GET")
    with urllib.request.urlopen(request, timeout=60) as response:
        return json.loads(response.read().decode("utf-8"))

def fetch_unsplash_photo(query: str, access_key: str) -> dict[str, str] | None:
    params = urllib.parse.urlencode({"query": query, "orientation": "landscape", "content_filter": "high"})
    headers = {"Authorization": f"Client-ID {access_key}", "Accept-Version": "v1"}
    try:
        photo = fetch_json(f"{UNSPLASH_RANDOM_URL}?{params}", headers)
    except urllib.error.HTTPError as exc:
        details = exc.read().decode("utf-8", errors="ignore")
        raise RuntimeError(f"Unsplash request failed ({exc.code}). {details}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"Network error while contacting Unsplash: {exc.reason}") from exc
    urls = photo.get("urls", {})
    user = photo.get("user", {})
    links = photo.get("links", {})
    profile = user.get("links", {}).get("html")
    return {
        "image_url": urls.get("regular", ""),
        "download_location": links.get("download_location", ""),
        "photographer": user.get("name", "Unsplash photographer"),
        "profile_url": with_referral(profile),
        "photo_url": with_referral(links.get("html")),
    }

def trigger_unsplash_download(download_location: str, access_key: str) -> None:
    if not download_location:
        return
    headers = {"Authorization": f"Client-ID {access_key}", "Accept-Version": "v1"}
    try:
        fetch_json(download_location, headers)
    except Exception:
        return

def fetch_image_bytes(url: str) -> bytes:
    if not url:
        return b""
    request = urllib.request.Request(url, headers={"User-Agent": APP_NAME}, method="GET")
    with urllib.request.urlopen(request, timeout=60) as response:
        return response.read()
def normalize_deck(data: dict[str, Any], topic: str, slide_count: int, theme_name: str, project_brief: list[str]) -> dict[str, Any]:
    slides: list[dict[str, Any]] = []
    for item in data.get("slides", []):
        bullets = [str(bullet).strip() for bullet in item.get("bullets", []) if str(bullet).strip()]
        if not bullets:
            continue
        slides.append({
            "title": str(item.get("title") or "Untitled Slide").strip(),
            "bullets": bullets[:4],
            "speaker_notes": str(item.get("speaker_notes") or "").strip(),
            "image": None,
        })
    while len(slides) < slide_count:
        fallback_bullets = [
            f"Expand on an important angle of {topic}",
            "Add a concrete example or supporting detail",
            "Close with a takeaway for the audience",
        ]
        if project_brief:
            fallback_bullets[0] = project_brief[min(len(slides), len(project_brief) - 1)]
        slides.append({
            "title": f"Key Point {len(slides) + 1}",
            "bullets": fallback_bullets,
            "speaker_notes": "Use this slide as a fallback if the model returns fewer slides than requested.",
            "image": None,
        })
    return {
        "title": str(data.get("title") or topic.title()).strip(),
        "subtitle": str(data.get("subtitle") or f"{theme_name} presentation on {topic}").strip(),
        "theme_name": theme_name,
        "slides": slides[:slide_count],
        "closing_message": str(data.get("closing_message") or "Questions and discussion").strip(),
        "project_brief": project_brief,
        "title_image": None,
    }

def enrich_deck_with_unsplash(deck: dict[str, Any], topic: str, access_key: str | None, include_title_image: bool = False) -> dict[str, Any]:
    # Attach one image candidate per slide so preview and export stay in sync.
    if not access_key:
        return deck
    if include_title_image:
        try:
            deck["title_image"] = fetch_unsplash_photo(topic, access_key)
        except Exception:
            deck["title_image"] = None
    for slide in deck["slides"]:
        query = f"{topic} {slide['title']}"
        try:
            slide["image"] = fetch_unsplash_photo(query, access_key)
        except Exception:
            slide["image"] = None
    return deck

def hex_to_rgb(value: str) -> RGBColor:
    cleaned = value.lstrip("#")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))

def apply_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)

def add_footer_text(slide, text: str, color: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.25))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = hex_to_rgb(color)

def get_image_geometry(layout_mode: str, slide_index: int) -> dict[str, float | str]:
    # Resolve whether the image sits left, right, or alternates across the deck.
    if layout_mode == "alternate":
        image_side = "left" if slide_index % 2 == 0 else "right"
    else:
        image_side = layout_mode
    if image_side == "left":
        return {"image_left": 0.7, "text_left": 6.7, "image_side": image_side}
    return {"image_left": 6.7, "text_left": 0.7, "image_side": image_side}

def build_presentation(deck: dict[str, Any], unsplash_key: str | None, theme: dict[str, str], layout_mode: str, presenter_name: str = "") -> bytes:
    # Export the final .pptx using the selected theme preset and image-placement mode.
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Install dependencies from requirements.txt first.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    primary = theme["primary"]
    secondary = theme["secondary"]
    accent = theme["accent"]
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(title_slide, primary)

    accent_band = title_slide.shapes.add_shape(1, Inches(0.95), Inches(0.95), Inches(1.6), Inches(0.16))
    accent_band.fill.solid()
    accent_band.fill.fore_color.rgb = hex_to_rgb(accent)
    accent_band.line.fill.background()

    title_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.2), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = deck["title"]
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = hex_to_rgb("#FFFFFF")
    title_paragraph.alignment = 1

    subtitle_box = title_slide.shapes.add_textbox(Inches(1.25), Inches(3.45), Inches(10.8), Inches(1.0))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = deck["subtitle"]
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(17)
    subtitle_paragraph.font.color.rgb = hex_to_rgb("#F6F0E8")
    subtitle_paragraph.alignment = 1

    if presenter_name.strip():
        presenter_box = title_slide.shapes.add_textbox(Inches(1.3), Inches(5.35), Inches(10.7), Inches(0.55))
        presenter_frame = presenter_box.text_frame
        presenter_frame.text = presenter_name.strip()
        presenter_paragraph = presenter_frame.paragraphs[0]
        presenter_paragraph.font.size = Pt(13)
        presenter_paragraph.font.bold = True
        presenter_paragraph.font.color.rgb = hex_to_rgb("#FFFFFF")
        presenter_paragraph.alignment = 1
    for index, slide_data in enumerate(deck["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide, secondary)
        geometry = get_image_geometry(layout_mode, index)
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.65))
        title_paragraph = title_box.text_frame.paragraphs[0]
        title_paragraph.text = slide_data["title"]
        title_paragraph.font.size = Pt(24)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = hex_to_rgb(primary)
        body_box = slide.shapes.add_textbox(Inches(float(geometry["text_left"])), Inches(1.35), Inches(5.2), Inches(4.7))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        for bullet_index, bullet in enumerate(slide_data["bullets"]):
            paragraph = body_frame.paragraphs[0] if bullet_index == 0 else body_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = hex_to_rgb(theme["ink"])
            paragraph.space_after = Pt(8)
        image = slide_data.get("image")
        if image and image.get("image_url"):
            image_bytes = fetch_image_bytes(image["image_url"])
            if image_bytes:
                slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(float(geometry["image_left"])), Inches(1.35), width=Inches(5.2), height=Inches(3.9))
                add_footer_text(slide, f"Photo: {image['photographer']} on Unsplash", primary)
                if unsplash_key:
                    trigger_unsplash_download(image.get("download_location", ""), unsplash_key)
        notes_text_frame = slide.notes_slide.notes_text_frame
        notes_text_frame.text = slide_data["speaker_notes"]
        if image and image.get("photographer"):
            notes_text_frame.text += (
                f"\n\nImage credit: {image['photographer']} | Profile: {image.get('profile_url', '')} | Photo: {image.get('photo_url', '')}"
            )
        accent_shape = slide.shapes.add_shape(1, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.14))
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = hex_to_rgb(accent)
        accent_shape.line.fill.background()
    closing_slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_background(closing_slide, primary)
    closing_title = closing_slide.shapes.title
    title_frame = closing_title.text_frame if closing_title else closing_slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(1.2)).text_frame
    title_frame.text = "Thank You"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = hex_to_rgb("#FFFFFF")
    message_box = closing_slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.2), Inches(1.4))
    message_frame = message_box.text_frame
    message_frame.text = deck["closing_message"]
    message_paragraph = message_frame.paragraphs[0]
    message_paragraph.font.size = Pt(20)
    message_paragraph.font.color.rgb = hex_to_rgb("#F6F0E8")
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

def make_file_stem(topic: str) -> str:
    safe_topic = re.sub(r"[^A-Za-z0-9_-]+", "-", topic.strip().lower()).strip("-")
    return safe_topic or "presentation"

def save_presentation_file(pptx_bytes: bytes, topic: str) -> Path:
    OUTPUT_DIR.mkdir(exist_ok=True)
    file_path = OUTPUT_DIR / f"{make_file_stem(topic)}.pptx"
    file_path.write_bytes(pptx_bytes)
    return file_path

api_key = get_api_key()
unsplash_key = get_unsplash_key()

with st.sidebar:
    st.markdown("### Presentation Builder")
    topic = st.text_input("Topic", placeholder="e.g. AI in healthcare")
    audience = st.text_input("Audience", value="Business stakeholders")
    tone = st.selectbox("Tone", ["Professional", "Educational", "Persuasive", "Executive summary"], index=0)
    slide_count = st.slider("Content slides", min_value=4, max_value=10, value=6)
    theme_name = st.selectbox("Background theme", list(THEME_PRESETS.keys()), index=1)
    image_layout_label = st.selectbox("Image placement", list(IMAGE_LAYOUT_OPTIONS.keys()), index=2)
    model = st.selectbox("Gemini model", [DEFAULT_MODEL, "gemini-2.0-flash", "gemini-1.5-flash"], index=0)
    use_unsplash = st.toggle("Include Unsplash images", value=bool(unsplash_key))
    presenter_name = st.text_input("Optional presenter name", placeholder="e.g. Ayo Ajayi")
    project_brief_text = st.text_area(
        "Optional project description",
        placeholder="Add up to five bullets, for example:\n- Problem we solve\n- Target users\n- Key differentiator\n- Business impact\n- Desired takeaway",
        height=160,
    )
    generate = st.button("Generate presentation", use_container_width=True)

selected_theme = THEME_PRESETS[theme_name]
layout_mode = IMAGE_LAYOUT_OPTIONS[image_layout_label]
project_brief = parse_project_brief(project_brief_text)
st.markdown(
    f"""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Sora:wght@500;600;700;800&family=Manrope:wght@400;500;600;700&display=swap');
        :root {{
            --bg: {selected_theme['surface']};
            --ink: {selected_theme['ink']};
            --muted: {selected_theme['muted']};
            --primary: {selected_theme['primary']};
            --accent: {selected_theme['accent']};
            --line: rgba(24, 26, 31, 0.10);
            --good-bg: rgba(61, 191, 161, 0.15);
            --good-ink: #1e6d58;
            --warn-bg: rgba(217, 123, 78, 0.14);
            --warn-ink: #9d542f;
        }}
        .stApp {{
            background:
                radial-gradient(circle at 12% 18%, rgba(217, 123, 78, 0.16), transparent 24%),
                radial-gradient(circle at 85% 14%, rgba(60, 191, 161, 0.12), transparent 20%),
                linear-gradient(145deg, #fbfaf7 0%, var(--bg) 52%, #eef2ee 100%);
            color: var(--ink);
        }}
        .block-container {{ max-width: 1220px; padding-top: 1.9rem; padding-bottom: 2rem; }}
        h1, h2, h3, h4 {{ font-family: "Sora", sans-serif; color: var(--ink); }}
        p, li, div[data-testid="stMarkdownContainer"], label {{ font-family: "Manrope", sans-serif; }}
        .hero, .card, .suggestion {{
            background: linear-gradient(180deg, rgba(255,255,255,0.82), rgba(255,255,255,0.68));
            border: 1px solid var(--line);
            box-shadow: 0 18px 60px rgba(35, 40, 44, 0.08);
            border-radius: 28px;
        }}
        .hero {{ padding: 3rem 2.4rem 2.7rem 2.4rem; position: relative; overflow: hidden; text-align: center; }}
        .hero::before {{ content: ""; position: absolute; left: -34px; top: -34px; width: 200px; height: 200px; border-radius: 999px; background: radial-gradient(circle, rgba(60, 191, 161, 0.14), transparent 68%); }}
        .hero::after {{ content: ""; position: absolute; right: -28px; top: -28px; width: 180px; height: 180px; border-radius: 999px; background: radial-gradient(circle, rgba(217, 123, 78, 0.20), transparent 68%); }}
        .eyebrow {{ display:inline-block; padding:0.5rem 0.9rem; border-radius:999px; background:rgba(255,255,255,0.84); border:1px solid rgba(24,26,31,0.08); font-size:0.82rem; font-weight:800; text-transform:uppercase; letter-spacing:0.12em; color:var(--primary); margin:0 auto; text-align:center; }}
        .hero-title {{ font-size: clamp(3.1rem, 6vw, 6rem); line-height: 0.94; letter-spacing: -0.05em; margin: 1rem auto 0.7rem auto; max-width: 11ch; text-wrap: balance; text-align:center; }}
        .hero-copy {{ color: var(--muted); max-width: 52rem; font-size: 1.1rem; line-height: 1.6; margin: 0 auto 1rem auto; text-wrap: pretty; text-align:center; }}
        .pill-row {{ display:flex; flex-wrap:wrap; gap:0.65rem; margin-top:1rem; justify-content:center; }}
        .pill, .status-pill {{ display:inline-block; padding:0.62rem 0.92rem; border-radius:999px; font-weight:700; }}
        .pill {{ background: rgba(24,26,31,0.05); color: var(--ink); }}
        .status-good {{ background: var(--good-bg); color: var(--good-ink); }}
        .status-warn {{ background: var(--warn-bg); color: var(--warn-ink); }}
        .section-grid {{ display:grid; grid-template-columns: repeat(3, minmax(0, 1fr)); gap:1rem; margin-top:1rem; }}
        .suggestion {{ padding:1rem 1.05rem; }}
        .suggestion h4 {{ margin:0 0 0.35rem 0; font-size:1.05rem; }}
        .suggestion p {{ margin:0; color:var(--muted); line-height:1.45; }}
        .card {{ padding:1.25rem 1.3rem; margin-top:1rem; }}
        .meta-grid {{ display:grid; grid-template-columns: repeat(4, minmax(0, 1fr)); gap:0.8rem; margin-top:0.85rem; }}
        .meta-item {{ border-radius:18px; padding:0.9rem 1rem; background:rgba(255,255,255,0.7); border:1px solid rgba(24,26,31,0.08); }}
        .meta-label {{ color:var(--muted); font-size:0.86rem; text-transform:uppercase; letter-spacing:0.04em; margin-bottom:0.25rem; }}
        .slide-card {{ background: rgba(255,255,255,0.78); border:1px solid rgba(24,26,31,0.08); border-radius:24px; padding:1rem 1.1rem; margin-bottom:1rem; }}
        .credit-text {{ color: var(--muted); font-size: 0.92rem; margin-top: 0.55rem; }}
        div[data-testid="stButton"] button, div[data-testid="stDownloadButton"] button {{ border-radius:16px; border:none; background:linear-gradient(135deg, var(--accent) 0%, var(--primary) 100%); color:white; font-weight:700; padding:0.75rem 1.1rem; }}
        @media (max-width: 900px) {{ .section-grid, .meta-grid {{ grid-template-columns:1fr; }} }}
    </style>
    """,
    unsafe_allow_html=True,
)

for key in ["deck", "pptx_bytes", "pptx_path", "last_error"]:
    if key not in st.session_state:
        st.session_state[key] = None

status_class = "status-good" if api_key else "status-warn"
status_text = "Gemini ready" if api_key else "Add GEMINI_API_KEY or GOOGLE_API_KEY"
image_status_class = "status-good" if unsplash_key else "status-warn"
image_status_text = "Unsplash ready" if unsplash_key else "Add UNSPLASH_ACCESS_KEY for images"

st.markdown(
    f"""
    <section class="hero">
        <span class="eyebrow">Group X</span>
        <h1 class="hero-title">  Bring your presentation ideas to life.</h1>
        <p class="hero-copy">
            Design sharper, more persuasive decks with a modern visual system built for polished first impressions.
        </p>
        <div class="pill-row">
            <span class="pill">{theme_name}</span>
            <span class="pill">{image_layout_label}</span>
            <span class="pill">{slide_count} content slides</span>
            <span class="status-pill {status_class}">{status_text}</span>
            <span class="status-pill {image_status_class}">{image_status_text}</span>
        </div>
    </section>
    """,
    unsafe_allow_html=True,
)

st.markdown(
    """
    <div class="section-grid">
        <div class="suggestion"><h4>Suggestion 1</h4><p>Use <strong>Midnight Luxe</strong> with <strong>Alternating left and right</strong> for investor or demo-day decks.</p></div>
        <div class="suggestion"><h4>Suggestion 2</h4><p>Use <strong>Neo Mint</strong> with <strong>Right aligned throughout</strong> for AI, product, and startup topics.</p></div>
        <div class="suggestion"><h4>Suggestion 3</h4><p>Use <strong>Sage Minimal</strong> plus a short five-bullet brief for calmer educational storytelling.</p></div>
    </div>
    """,
    unsafe_allow_html=True,
)

if generate:
    # Clear previous output so each run shows only the latest presentation.
    st.session_state.last_error = None
    st.session_state.deck = None
    st.session_state.pptx_bytes = None
    st.session_state.pptx_path = None
    if not topic.strip():
        st.session_state.last_error = "Enter a presentation topic first."
    elif not api_key:
        st.session_state.last_error = "No Gemini API key found. Add GEMINI_API_KEY or GOOGLE_API_KEY to your .env file."
    elif use_unsplash and not unsplash_key:
        st.session_state.last_error = "Unsplash images are enabled, but UNSPLASH_ACCESS_KEY is missing from .env."
    else:
        prompt = build_prompt(topic.strip(), audience.strip() or "General audience", tone, slide_count, theme_name, project_brief)
        with st.spinner("Generating presentation..."):
            try:
                raw_deck = call_gemini(prompt, api_key, model)
                deck = normalize_deck(raw_deck, topic.strip(), slide_count, theme_name, project_brief)
                deck = enrich_deck_with_unsplash(
                    deck,
                    topic.strip(),
                    unsplash_key if use_unsplash else None,
                    include_title_image=False,
                )
                pptx_bytes = build_presentation(
                    deck,
                    unsplash_key if use_unsplash else None,
                    selected_theme,
                    layout_mode,
                    presenter_name=presenter_name,
                )
                pptx_path = save_presentation_file(pptx_bytes, topic.strip())
            except Exception as exc:
                st.session_state.last_error = str(exc)
            else:
                st.session_state.deck = deck
                st.session_state.pptx_bytes = pptx_bytes
                st.session_state.pptx_path = str(pptx_path)

if st.session_state.last_error:
    st.error(st.session_state.last_error)

left_col, right_col = st.columns([1.35, 0.85], gap="large")
with left_col:
    st.markdown(
        f"""
        <div class="card">
            <h3 style="margin-top:0;">Current setup</h3>
            <div class="meta-grid">
                <div class="meta-item"><div class="meta-label">Theme</div><div><strong>{theme_name}</strong></div></div>
                <div class="meta-item"><div class="meta-label">Image layout</div><div><strong>{image_layout_label}</strong></div></div>
                <div class="meta-item"><div class="meta-label">Optional brief</div><div><strong>{len(project_brief)}/5 bullets</strong></div></div>
                <div class="meta-item"><div class="meta-label">Presenter</div><div><strong>{presenter_name if presenter_name.strip() else "Not set"}</strong></div></div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    deck = st.session_state.deck
    if deck:
        tabs = st.tabs(["Slide Preview", "Project Brief", "JSON Output"])
        with tabs[0]:
            st.subheader(deck["title"])
            st.caption(deck["subtitle"])
            for index, slide in enumerate(deck["slides"], start=1):
                st.markdown(f"<div class='slide-card'><strong>Slide {index}: {slide['title']}</strong></div>", unsafe_allow_html=True)
                geometry = get_image_geometry(layout_mode, index - 1)
                image = slide.get("image")
                if image and image.get("image_url"):
                    if geometry["image_side"] == "left":
                        image_col, text_col = st.columns([0.95, 1.05], gap="medium")
                    else:
                        text_col, image_col = st.columns([1.05, 0.95], gap="medium")
                else:
                    text_col = st.container()
                    image_col = None
                with text_col:
                    st.markdown("\n".join(f"- {bullet}" for bullet in slide["bullets"]))
                    st.caption(f"Speaker notes: {slide['speaker_notes']}")
                if image_col and image:
                    with image_col:
                        st.image(image["image_url"], use_container_width=True)
                        photographer = image.get("photographer", "Unsplash photographer")
                        photo_url = image.get("photo_url", "")
                        profile_url = image.get("profile_url", "")
                        st.markdown(
                            f"<p class='credit-text'>Photo by <a href='{profile_url}' target='_blank'>{photographer}</a> on <a href='{photo_url}' target='_blank'>Unsplash</a></p>",
                            unsafe_allow_html=True,
                        )
            st.markdown(f"<div class='slide-card'><strong>Final slide</strong><p>{deck['closing_message']}</p></div>", unsafe_allow_html=True)
        with tabs[1]:
            if deck["project_brief"]:
                st.markdown("\n".join(f"- {item}" for item in deck["project_brief"]))
            else:
                st.info("No project description bullets were provided for this run.")
        with tabs[2]:
            st.json(deck)
    else:
        st.markdown(
            """
            <div class="card">
                <h3 style="margin-top:0;">Ready to generate</h3>
                <p style="color:var(--muted);margin-bottom:0;">Choose your theme, image placement strategy, and optional five-bullet brief, then generate the deck.</p>
            </div>
            """,
            unsafe_allow_html=True,
        )

with right_col:
    st.markdown(
        f"""
        <div class="card">
            <h3 style="margin-top:0;">Theme preview</h3>
            <p style="color:var(--muted);">{selected_theme['description']}</p>
            <div class="pill-row">
                <span class="pill">Primary {selected_theme['primary']}</span>
                <span class="pill">Surface {selected_theme['surface']}</span>
                <span class="pill">Accent {selected_theme['accent']}</span>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    if st.session_state.pptx_bytes:
        safe_topic = make_file_stem(topic)
        st.download_button(
            "Download PowerPoint",
            data=st.session_state.pptx_bytes,
            file_name=f"{safe_topic}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
        st.success("Presentation file created.")
        if st.session_state.pptx_path:
            st.code(st.session_state.pptx_path, language="text")
    else:
        st.markdown(
            """
            <div class="card">
                <h3 style="margin-top:0;">Output</h3>
                <p style="color:var(--muted);margin-bottom:0;">Your `.pptx` file will be saved in `generated_presentations/` and available for download here.</p>
            </div>
            """,
            unsafe_allow_html=True,
        )
    if Presentation is None:
        st.error("python-pptx is not installed in this environment yet. Install requirements before exporting decks.")

st.caption("Gemini generates the slide structure, Unsplash provides optional images, and the app exports a themed PowerPoint with configurable image placement.")
